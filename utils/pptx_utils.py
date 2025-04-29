from pptx import Presentation
from pathlib import Path
from config import OUTPUT_DIR
import win32com.client
import pythoncom
import os

def handle_pptx(file):
    try:
        # Initialize COM in this thread
        pythoncom.CoInitialize()
        
        # First extract text using python-pptx
        prs = Presentation(file)
        text_content = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        
        text = "\n".join(text_content)
        
        # Then convert to PDF using PowerPoint
        powerpoint = win32com.client.Dispatch("Powerpoint.Application")
        powerpoint.Visible = False
        
        pdf_path = OUTPUT_DIR / file.with_suffix(".pdf").name
        deck = powerpoint.Presentations.Open(str(file.absolute()))
        deck.SaveAs(str(pdf_path.absolute()), 32)  # 32 is the PDF format code
        deck.Close()
        powerpoint.Quit()
        
        return pdf_path, text
        
    except Exception as e:
        print(f"Error processing PowerPoint file {file}: {str(e)}")
        return None, f"Error: {str(e)}"
    finally:
        pythoncom.CoUninitialize() 